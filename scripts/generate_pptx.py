#!/usr/bin/env python3
"""
Generate weekly Coach Playbook PowerPoint decks from public/training-plans.json.

One .pptx per week, written to public/plans/<weekOf>.pptx.
Each deck has a polished cover, two section dividers, and one drill slide per drill
with a real soccer-field diagram, coaching cues, and equipment setup.

Run from repo root:  python3 scripts/generate_pptx.py
"""

from __future__ import annotations

import json
import math
import os
import sys
from datetime import datetime

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
JSON_PATH = os.path.join(ROOT, "public", "training-plans.json")
OUT_DIR = os.path.join(ROOT, "public", "plans")

# ── Brand palette ────────────────────────────────────────────────────────────
GREEN = RGBColor(0x22, 0xA2, 0x4D)
GREEN_BRIGHT = RGBColor(0x4A, 0xDE, 0x80)
GREEN_DEEP = RGBColor(0x14, 0x5A, 0x2A)
FIELD_LIGHT = RGBColor(0x3E, 0x96, 0x4F)
FIELD_DARK = RGBColor(0x2A, 0x70, 0x37)
BG_DARK = RGBColor(0x08, 0x0D, 0x0A)
BG_PANEL = RGBColor(0x10, 0x18, 0x12)
BG_CARD = RGBColor(0x18, 0x22, 0x1B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_SOFT = RGBColor(0xE8, 0xED, 0xE9)
TEXT_MID = RGBColor(0xB4, 0xC3, 0xB8)
TEXT_DIM = RGBColor(0x7B, 0x8A, 0x7F)
BORDER = RGBColor(0x1F, 0x2D, 0x23)

GOLD = RGBColor(0xF5, 0xA6, 0x23)
ORANGE = RGBColor(0xEE, 0x7B, 0x2C)
BLUE = RGBColor(0x4A, 0x90, 0xE2)
PURPLE = RGBColor(0x7C, 0x5D, 0xFA)
RED = RGBColor(0xE2, 0x4A, 0x4A)
PINK_RED = RGBColor(0xFF, 0x6B, 0x6B)
CYAN = RGBColor(0x4F, 0xD1, 0xC5)

TAG_COLOR = {
    "Opening Game": GOLD,
    "Lesson 1": BLUE,
    "Lesson 2": PURPLE,
    "Scrimmage": GREEN_BRIGHT,
}

MONTHS_FULL = ["January", "February", "March", "April", "May", "June",
               "July", "August", "September", "October", "November", "December"]


# ── Low-level XML helpers ────────────────────────────────────────────────────

def _qn(tag: str) -> str:
    return qn(tag)


def add_outer_shadow(shape, *, blur=60000, distance=38100, alpha=55000):
    """Add a soft drop shadow to any shape (works on shapes with .fill)."""
    spPr = shape._element.spPr
    # Remove any existing effect list
    for old in spPr.findall(_qn("a:effectLst")):
        spPr.remove(old)
    effectLst = etree.SubElement(spPr, _qn("a:effectLst"))
    shdw = etree.SubElement(effectLst, _qn("a:outerShdw"))
    shdw.set("blurRad", str(blur))
    shdw.set("dist", str(distance))
    shdw.set("dir", "5400000")  # 90° down
    shdw.set("algn", "ctr")
    shdw.set("rotWithShape", "0")
    clr = etree.SubElement(shdw, _qn("a:srgbClr"))
    clr.set("val", "000000")
    alpha_el = etree.SubElement(clr, _qn("a:alpha"))
    alpha_el.set("val", str(alpha))


def set_line_arrow(line, *, head=None, tail="arrow", head_size="med", tail_size="med"):
    """Add arrowheads to a line/connector. head/tail are types like 'arrow', 'triangle', 'oval'."""
    ln = line.line._get_or_add_ln()
    # Remove existing arrows
    for el in ln.findall(_qn("a:headEnd")):
        ln.remove(el)
    for el in ln.findall(_qn("a:tailEnd")):
        ln.remove(el)
    if head:
        he = etree.SubElement(ln, _qn("a:headEnd"))
        he.set("type", head)
        he.set("w", head_size)
        he.set("len", head_size)
    if tail:
        te = etree.SubElement(ln, _qn("a:tailEnd"))
        te.set("type", tail)
        te.set("w", tail_size)
        te.set("len", tail_size)


def set_line_dash(line, dash_style: str):
    """dash_style: 'solid', 'dash', 'dashDot', 'lgDash', 'sysDash', 'sysDot'."""
    ln = line.line._get_or_add_ln()
    for old in ln.findall(_qn("a:prstDash")):
        ln.remove(old)
    prst = etree.SubElement(ln, _qn("a:prstDash"))
    prst.set("val", dash_style)


def fill_gradient(shape, color_top: RGBColor, color_bottom: RGBColor, angle_deg: float = 90.0):
    """Apply a 2-stop linear gradient fill."""
    spPr = shape._element.spPr
    # Remove existing fill elements
    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill", "a:pattFill"):
        for old in spPr.findall(_qn(tag)):
            spPr.remove(old)
    grad = etree.SubElement(spPr, _qn("a:gradFill"))
    grad.set("flip", "none")
    grad.set("rotWithShape", "1")
    gsLst = etree.SubElement(grad, _qn("a:gsLst"))
    gs0 = etree.SubElement(gsLst, _qn("a:gs")); gs0.set("pos", "0")
    c0 = etree.SubElement(gs0, _qn("a:srgbClr")); c0.set("val", "{:02X}{:02X}{:02X}".format(color_top[0], color_top[1], color_top[2]))
    gs1 = etree.SubElement(gsLst, _qn("a:gs")); gs1.set("pos", "100000")
    c1 = etree.SubElement(gs1, _qn("a:srgbClr")); c1.set("val", "{:02X}{:02X}{:02X}".format(color_bottom[0], color_bottom[1], color_bottom[2]))
    lin = etree.SubElement(grad, _qn("a:lin"))
    lin.set("ang", str(int(angle_deg * 60000)))
    lin.set("scaled", "0")
    # Reorder: spPr children expect fill before line; insert near top
    spPr.insert(1, grad)
    # Already added; second insert is harmless since we removed all fills above


# ── Shape builders ───────────────────────────────────────────────────────────

def set_solid_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def set_line(shape, color: RGBColor, width_pt: float):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def rect(slide, x, y, w, h, *, fill=None, gradient=None, line=None, line_w=1.0, rounded=False, shadow=False):
    shape_id = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_id, x, y, w, h)
    if rounded:
        try:
            s.adjustments[0] = 0.12
        except Exception:
            pass
    if gradient:
        # gradient = (top_color, bottom_color, angle)
        fill_gradient(s, *gradient)
    elif fill is not None:
        set_solid_fill(s, fill)
    else:
        s.fill.background()
    if line is None:
        no_line(s)
    else:
        set_line(s, line, line_w)
    if shadow:
        add_outer_shadow(s)
    return s


def oval(slide, x, y, w, h, *, fill, line=None, line_w=1.0, shadow=False):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    set_solid_fill(s, fill)
    if line is None:
        no_line(s)
    else:
        set_line(s, line, line_w)
    if shadow:
        add_outer_shadow(s)
    return s


def triangle(slide, cx, cy, size, *, fill, rotation=0):
    s = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, cx - size / 2, cy - size / 2, size, size)
    set_solid_fill(s, fill)
    no_line(s)
    if rotation:
        s.rotation = rotation
    return s


def textbox(slide, x, y, w, h, text, *, size=14, bold=False, color=WHITE,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
            italic=False, letter_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    lines = text if isinstance(text, list) else text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        if letter_spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(letter_spacing))
    return tb


def connector(slide, x1, y1, x2, y2, *, color=WHITE, width=2.0,
              dash=None, arrow_tail="arrow", arrow_head=None,
              tail_size="med"):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if dash:
        set_line_dash(line, dash)
    set_line_arrow(line, head=arrow_head, tail=arrow_tail, tail_size=tail_size)
    return line


def chip(slide, x, y, label, *, bg, fg=WHITE, size=10, padding_x=Inches(0.15), height=Inches(0.3)):
    """A rounded chip/pill with text."""
    # Estimate width from label length
    char_w = size * 7.4 / 12  # rough EMU per char
    width_in = max(0.5, (len(label) * char_w / 72) + 0.4)
    w = Inches(width_in)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, height)
    try:
        s.adjustments[0] = 0.5
    except Exception:
        pass
    set_solid_fill(s, bg)
    no_line(s)
    textbox(slide, x, y, w, height, label,
            size=size, bold=True, color=fg, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            letter_spacing=80)
    return s, w


# ── Soccer-field components ──────────────────────────────────────────────────

def soccer_field(slide, x, y, w, h, *, halves=False, penalty_boxes=False,
                 center_circle=False, dark=False):
    """Draw a soccer field with markings. Returns the field rect bounds."""
    field_top = FIELD_DARK if dark else FIELD_LIGHT
    field_bot = GREEN_DEEP if dark else FIELD_DARK
    base = rect(slide, x, y, w, h, gradient=(field_top, field_bot, 135.0))
    no_line(base)

    # Outer boundary
    inset = Emu(int(min(w, h) * 0.03))
    boundary = rect(slide, x + inset, y + inset, w - inset * 2, h - inset * 2,
                    line=WHITE, line_w=2.5)
    boundary.fill.background()

    if halves:
        connector(slide, x + w / 2, y + inset, x + w / 2, y + h - inset,
                  color=WHITE, width=1.5, dash="dash", arrow_tail=None)

    if center_circle:
        cr = min(w, h) * 0.13
        circle = oval(slide, x + w / 2 - cr, y + h / 2 - cr, cr * 2, cr * 2,
                      fill=field_top, line=WHITE, line_w=1.5)
        # Hide fill
        circle.fill.background()
        circle.line.color.rgb = WHITE
        circle.line.width = Pt(1.5)
        # Center dot
        dot_r = Emu(int(min(w, h) * 0.012))
        oval(slide, x + w / 2 - dot_r, y + h / 2 - dot_r, dot_r * 2, dot_r * 2, fill=WHITE)

    if penalty_boxes:
        box_w = w * 0.18
        box_h = h * 0.45
        # Left box
        lb = rect(slide, x + inset, y + h / 2 - box_h / 2, box_w, box_h, line=WHITE, line_w=1.5)
        lb.fill.background()
        # Right box
        rb = rect(slide, x + w - inset - box_w, y + h / 2 - box_h / 2, box_w, box_h, line=WHITE, line_w=1.5)
        rb.fill.background()
        # Six-yard boxes
        sw = w * 0.07
        sh = h * 0.22
        ls = rect(slide, x + inset, y + h / 2 - sh / 2, sw, sh, line=WHITE, line_w=1.5)
        ls.fill.background()
        rs = rect(slide, x + w - inset - sw, y + h / 2 - sh / 2, sw, sh, line=WHITE, line_w=1.5)
        rs.fill.background()

    return x, y, w, h


def grid_field(slide, x, y, w, h, *, label_size=None):
    """Small training grid with hash boundary."""
    top = FIELD_LIGHT
    bot = FIELD_DARK
    base = rect(slide, x, y, w, h, gradient=(top, bot, 135.0))
    inset = Emu(int(min(w, h) * 0.04))
    boundary = rect(slide, x + inset, y + inset, w - inset * 2, h - inset * 2,
                    line=WHITE, line_w=2.0)
    boundary.fill.background()
    set_line_dash(boundary, "lgDash")


def goal(slide, cx, cy, *, width_in=0.8, depth_in=0.18, orientation="horizontal"):
    """Soccer goal — white frame, mouth facing into play."""
    if orientation == "horizontal":
        w = Inches(width_in)
        d = Inches(depth_in)
        x = cx - w / 2
        y = cy - d / 2
        # Goal net background (translucent)
        net = rect(slide, x, y, w, d, fill=BG_PANEL)
        net.fill.transparency = 0.5  # not supported reliably; ok if ignored
        # Posts
        post_w = Pt(3)
        rect(slide, x, y, Inches(0.04), d, fill=WHITE)
        rect(slide, x + w - Inches(0.04), y, Inches(0.04), d, fill=WHITE)
        rect(slide, x, y, w, Inches(0.04), fill=WHITE)
    else:
        d = Inches(width_in)
        w = Inches(depth_in)
        x = cx - w / 2
        y = cy - d / 2
        net = rect(slide, x, y, w, d, fill=BG_PANEL)
        rect(slide, x, y, Inches(0.04), d, fill=WHITE)
        rect(slide, x + w - Inches(0.04), y, Inches(0.04), d, fill=WHITE)
        rect(slide, x, y, w, Inches(0.04), fill=WHITE)


def player(slide, cx, cy, team_color, label="", *, radius_in=0.16, ring=WHITE):
    r = Inches(radius_in)
    oval(slide, cx - r, cy - r, r * 2, r * 2, fill=team_color, line=ring, line_w=1.5, shadow=True)
    if label:
        textbox(slide, cx - r, cy - Inches(radius_in * 0.7), r * 2, Inches(radius_in * 1.4),
                label, size=int(radius_in * 70), bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def ball(slide, cx, cy, *, radius_in=0.085):
    r = Inches(radius_in)
    oval(slide, cx - r, cy - r, r * 2, r * 2, fill=WHITE, line=BG_DARK, line_w=1.0, shadow=True)
    # Pentagon hint
    inner_r = Inches(radius_in * 0.45)
    oval(slide, cx - inner_r, cy - inner_r, inner_r * 2, inner_r * 2, fill=RGBColor(0x33, 0x33, 0x33))


def cone(slide, cx, cy, *, size_in=0.22):
    """Orange traffic cone — triangle with shadow."""
    size = Inches(size_in)
    base_h = Inches(size_in * 0.18)
    # Base
    rect(slide, cx - size / 2, cy + size / 2 - base_h / 2, size, base_h,
         fill=RGBColor(0xC9, 0x5A, 0x14), shadow=True)
    # Cone body
    t = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, cx - size / 2, cy - size / 2, size, size)
    fill_gradient(t, GOLD, ORANGE, 90.0)
    no_line(t)
    # Reflective band (lighter stripe)
    band_h = Inches(size_in * 0.12)
    rect(slide, cx - size / 3, cy - size / 8, size * 2 / 3, band_h,
         fill=RGBColor(0xFF, 0xE5, 0xB8))


def pass_arrow(slide, x1, y1, x2, y2, *, color=None):
    """Dashed pass arrow with arrowhead."""
    connector(slide, x1, y1, x2, y2,
              color=color or WHITE, width=2.0, dash="dash",
              arrow_tail="arrow", tail_size="med")


def run_arrow(slide, x1, y1, x2, y2, *, color=None):
    """Solid run arrow."""
    connector(slide, x1, y1, x2, y2,
              color=color or GOLD, width=2.5,
              arrow_tail="triangle", tail_size="lg")


def dribble_arrow(slide, x1, y1, x2, y2, *, color=None):
    """Wavy dribble arrow (simulated with dash-dot)."""
    connector(slide, x1, y1, x2, y2,
              color=color or WHITE, width=2.0, dash="dashDot",
              arrow_tail="triangle", tail_size="med")


def zone_box(slide, x, y, w, h, *, color):
    """A semi-transparent zone callout."""
    s = rect(slide, x, y, w, h, line=color, line_w=1.5, rounded=True)
    s.fill.background()
    set_line_dash(s, "dash")


# ── Diagram dispatcher ───────────────────────────────────────────────────────

def classify_diagram(title: str, tag: str) -> str:
    t = (title or "").lower()
    if any(k in t for k in ["sharks", "pass tag", "king of the hill", "3-team", "world cup knockout"]):
        return "grid_chase"
    if any(k in t for k in ["red light", "follow the coach"]):
        return "grid_listen"
    if any(k in t for k in ["cone path", "friendly foot", "close-control", "slalom", "beat the cone"]):
        return "slalom"
    if any(k in t for k in ["push & chase", "push and chase", "race to the cone"]):
        return "straight_lane"
    if "first touch & drive" in t or "drive & finish" in t or ("drive" in t and "finish" in t):
        return "drive_finish"
    if "first touch & drive" in t:
        return "first_touch_drive"
    if any(k in t for k in ["toe tap", "skills showcase"]):
        return "freeform_grid"
    if any(k in t for k in ["pass to the coach", "pass the coach"]):
        return "circle_pass"
    if any(k in t for k in ["stop & pass", "stop and pass", "pass through the gate", "2-touch", "partner passing"]):
        return "pairs_passing"
    if any(k in t for k in ["pass & move", "triangles"]):
        return "triangle"
    if any(k in t for k in ["give-and-go", "wall pass", "combination", "1-2"]):
        return "give_and_go"
    if any(k in t for k in ["2v1"]):
        return "two_v_one"
    if any(k in t for k in ["1v1"]):
        return "one_v_one"
    if any(k in t for k in ["go around the coach", "go around"]):
        return "around_coach"
    if any(k in t for k in ["two moves"]):
        return "moves_pair"
    if any(k in t for k in ["knock down", "shoot the ball", "score 5", "hot potato"]):
        return "shooting_line"
    if any(k in t for k in ["finishing from a pass"]):
        return "finishing_pass"
    if any(k in t for k in ["everyone scores", "mini game"]):
        return "mini_game"
    if any(k in t for k in ["small-sided scrimmage", "season finale", "scrimmage"]):
        return "scrimmage"
    if any(k in t for k in ["position", "free play"]):
        return "formation"
    return "freeform_grid"


def draw_diagram(slide, x, y, w, h, *, title: str, tag: str, space: str):
    """Draw a polished, drill-appropriate field diagram inside (x, y, w, h)."""
    kind = classify_diagram(title, tag)
    cx = x + w / 2
    cy = y + h / 2
    band = TAG_COLOR.get(tag, GREEN)

    if kind in ("grid_chase", "grid_listen", "freeform_grid"):
        grid_field(slide, x, y, w, h)

        if kind == "grid_chase":
            # Minnows (with balls) + sharks
            minnow_positions = [(0.18, 0.28), (0.32, 0.55), (0.55, 0.32),
                                 (0.72, 0.6), (0.85, 0.4), (0.28, 0.78), (0.6, 0.82)]
            for px, py in minnow_positions:
                mcx = x + w * px
                mcy = y + h * py
                player(slide, mcx, mcy, BLUE, "M")
                ball(slide, mcx + Inches(0.22), mcy + Inches(0.05))
            for px, py in [(0.45, 0.48), (0.7, 0.65)]:
                scx = x + w * px
                scy = y + h * py
                player(slide, scx, scy, RED, "S", radius_in=0.18)
                # Movement toward minnow
                run_arrow(slide, scx + Inches(0.15), scy, scx + Inches(0.55), scy - Inches(0.15), color=PINK_RED)

        elif kind == "grid_listen":
            # All players have balls, spread
            positions = [(0.18, 0.3), (0.34, 0.55), (0.55, 0.35), (0.72, 0.6),
                          (0.85, 0.4), (0.28, 0.78), (0.6, 0.8), (0.45, 0.22)]
            for i, (px, py) in enumerate(positions):
                pcx = x + w * px
                pcy = y + h * py
                player(slide, pcx, pcy, GOLD, str(i + 1))
                ball(slide, pcx + Inches(0.2), pcy)
            # Coach center
            player(slide, cx, cy, GREEN, "C", radius_in=0.22)

        else:  # freeform_grid (each has their own ball)
            cols, rows = 4, 2
            for cc in range(cols):
                for rr in range(rows):
                    px = 0.18 + cc * 0.21
                    py = 0.35 + rr * 0.32
                    pcx = x + w * px
                    pcy = y + h * py
                    player(slide, pcx, pcy, GOLD, "P")
                    ball(slide, pcx, pcy + Inches(0.22))

    elif kind == "slalom":
        grid_field(slide, x, y, w, h)
        # Cone slalom horizontal
        n = 6
        ys = [0.45, 0.55, 0.45, 0.55, 0.45, 0.55]
        positions = []
        for i in range(n):
            px = 0.2 + (0.6 / (n - 1)) * i
            py = ys[i]
            cone(slide, x + w * px, y + h * py)
            positions.append((x + w * px, y + h * py))

        # Curved dribble path implied with multiple short arrows
        start = (x + w * 0.1, y + h * 0.5)
        player(slide, start[0], start[1], GOLD, "P")
        ball(slide, start[0] + Inches(0.18), start[1])
        prev = (start[0] + Inches(0.2), start[1])
        for px, py in positions:
            # offset target slightly past the cone
            tgt = (px + Inches(0.08), py - Inches(0.15) if py > y + h * 0.5 else py + Inches(0.15))
            dribble_arrow(slide, prev[0], prev[1], tgt[0], tgt[1], color=GOLD)
            prev = tgt
        finish = (x + w * 0.92, y + h * 0.5)
        dribble_arrow(slide, prev[0], prev[1], finish[0], finish[1], color=GOLD)

    elif kind == "straight_lane":
        grid_field(slide, x, y, w, h)
        start = (x + w * 0.12, cy)
        end = (x + w * 0.88, cy)
        # Lane markers
        for i in range(1, 5):
            mx = x + w * (0.12 + (0.76 / 4) * i)
            connector(slide, mx, cy - Inches(0.05), mx, cy + Inches(0.05),
                      color=WHITE, width=1.0, arrow_tail=None)
        cone(slide, end[0], end[1])
        player(slide, start[0], start[1], GOLD, "P")
        ball(slide, start[0] + Inches(0.18), start[1])
        run_arrow(slide, start[0] + Inches(0.22), start[1],
                   end[0] - Inches(0.2), end[1], color=GOLD)

    elif kind == "drive_finish":
        soccer_field(slide, x, y, w, h, penalty_boxes=True)
        # Goal on right
        goal(slide, x + w - Inches(0.18), cy, width_in=1.2, orientation="vertical")
        # Player with ball mid-field
        pcx = x + w * 0.25
        player(slide, pcx, cy, GOLD, "P")
        ball(slide, pcx + Inches(0.2), cy)
        # Coach passes
        coach_x = x + w * 0.18
        coach_y = y + h * 0.25
        player(slide, coach_x, coach_y, GREEN, "C")
        pass_arrow(slide, coach_x + Inches(0.15), coach_y + Inches(0.05),
                    pcx - Inches(0.05), cy - Inches(0.1))
        # Drive + shot
        drive_to = (x + w * 0.55, cy)
        dribble_arrow(slide, pcx + Inches(0.2), cy, drive_to[0], drive_to[1], color=GOLD)
        shot_to = (x + w - Inches(0.25), cy - Inches(0.15))
        run_arrow(slide, drive_to[0] + Inches(0.05), drive_to[1],
                   shot_to[0], shot_to[1], color=RED)
        # Optional defender
        player(slide, x + w * 0.7, cy + Inches(0.3), PINK_RED, "D")

    elif kind == "first_touch_drive":
        soccer_field(slide, x, y, w, h)
        pcx = x + w * 0.3
        player(slide, pcx, cy, GOLD, "A")
        coach_x = x + w * 0.15
        player(slide, coach_x, cy, GREEN, "C")
        pass_arrow(slide, coach_x + Inches(0.12), cy, pcx - Inches(0.1), cy)
        drive_to_x = x + w * 0.75
        dribble_arrow(slide, pcx + Inches(0.12), cy, drive_to_x, cy, color=GOLD)
        ball(slide, pcx + Inches(0.15), cy)

    elif kind == "circle_pass":
        grid_field(slide, x, y, w, h)
        r_in = (min(w, h) / 3) / Inches(1)
        n = 7
        for i in range(n):
            ang = (2 * math.pi / n) * i - math.pi / 2
            pcx = cx + Inches(r_in) * math.cos(ang)
            pcy = cy + Inches(r_in) * math.sin(ang)
            player(slide, pcx, pcy, GOLD, "P")
            pass_arrow(slide, pcx - Inches(0.1) * math.cos(ang),
                        pcy - Inches(0.1) * math.sin(ang),
                        cx + Inches(0.18) * math.cos(ang),
                        cy + Inches(0.18) * math.sin(ang))
        player(slide, cx, cy, GREEN, "C", radius_in=0.22)
        ball(slide, cx + Inches(0.3), cy + Inches(0.05))

    elif kind == "pairs_passing":
        grid_field(slide, x, y, w, h)
        for row, py_frac in enumerate([0.32, 0.72]):
            yrow = y + h * py_frac
            a = (x + w * 0.18, yrow)
            b = (x + w * 0.82, yrow)
            cone(slide, x + w * 0.46, yrow - Inches(0.18))
            cone(slide, x + w * 0.46, yrow + Inches(0.18))
            player(slide, a[0], a[1], GOLD, "A")
            player(slide, b[0], b[1], BLUE, "B")
            ball(slide, a[0] + Inches(0.2), yrow)
            pass_arrow(slide, a[0] + Inches(0.18), yrow - Inches(0.04),
                        b[0] - Inches(0.18), yrow - Inches(0.04))
            pass_arrow(slide, b[0] - Inches(0.18), yrow + Inches(0.05),
                        a[0] + Inches(0.18), yrow + Inches(0.05), color=GREEN_BRIGHT)

    elif kind == "triangle":
        grid_field(slide, x, y, w, h)
        positions = [(0.26, 0.78), (0.5, 0.22), (0.74, 0.78)]
        labels = ["A", "B", "C"]
        colors = [GOLD, BLUE, PURPLE]
        pts = []
        for (px, py), lab, col in zip(positions, labels, colors):
            pcx = x + w * px
            pcy = y + h * py
            player(slide, pcx, pcy, col, lab)
            pts.append((pcx, pcy))
        for i in range(3):
            ax, ay = pts[i]
            bx, by = pts[(i + 1) % 3]
            pass_arrow(slide, ax + Inches(0.1), ay, bx - Inches(0.1), by)
        ball(slide, pts[0][0] + Inches(0.18), pts[0][1])

    elif kind == "give_and_go":
        soccer_field(slide, x, y, w, h)
        a = (x + w * 0.2, cy + Inches(0.35))
        b = (x + w * 0.55, cy - Inches(0.4))
        defender = (x + w * 0.45, cy)
        end = (x + w * 0.78, cy + Inches(0.1))
        player(slide, a[0], a[1], GOLD, "A")
        player(slide, b[0], b[1], BLUE, "B")
        player(slide, defender[0], defender[1], PINK_RED, "D")
        ball(slide, a[0] + Inches(0.18), a[1])
        # 1: A passes to B
        pass_arrow(slide, a[0] + Inches(0.15), a[1] - Inches(0.05),
                    b[0] - Inches(0.1), b[1] + Inches(0.08))
        # 2: A runs past defender
        run_arrow(slide, a[0] + Inches(0.1), a[1] + Inches(0.05),
                   end[0] - Inches(0.1), end[1] + Inches(0.05))
        # 3: B returns pass into A's path
        pass_arrow(slide, b[0], b[1] + Inches(0.15),
                    end[0] - Inches(0.05), end[1], color=GREEN_BRIGHT)

    elif kind == "two_v_one":
        soccer_field(slide, x, y, w, h, penalty_boxes=True)
        goal(slide, x + w - Inches(0.18), cy, width_in=1.2, orientation="vertical")
        a = (x + w * 0.2, cy + Inches(0.3))
        a2 = (x + w * 0.3, cy - Inches(0.4))
        d = (x + w * 0.55, cy)
        player(slide, a[0], a[1], GOLD, "A1")
        player(slide, a2[0], a2[1], GOLD, "A2")
        player(slide, d[0], d[1], PINK_RED, "D")
        ball(slide, a[0] + Inches(0.18), a[1])
        pass_arrow(slide, a[0] + Inches(0.15), a[1] - Inches(0.05),
                    a2[0] - Inches(0.05), a2[1] + Inches(0.1))
        run_arrow(slide, a[0] + Inches(0.05), a[1] - Inches(0.1),
                   x + w * 0.65, cy + Inches(0.2))
        run_arrow(slide, a2[0] + Inches(0.1), a2[1] + Inches(0.05),
                   x + w - Inches(0.3), cy - Inches(0.15), color=RED)

    elif kind in ("one_v_one", "around_coach", "moves_pair"):
        grid_field(slide, x, y, w, h)
        att = (x + w * 0.22, cy)
        defx = x + w * 0.55
        player(slide, att[0], att[1], GOLD, "A")
        ball(slide, att[0] + Inches(0.18), cy)
        player(slide, defx, cy, PINK_RED, "D")
        # Two move options
        run_arrow(slide, att[0] + Inches(0.15), cy - Inches(0.05),
                   x + w * 0.88, y + h * 0.28, color=GOLD)
        run_arrow(slide, att[0] + Inches(0.15), cy + Inches(0.05),
                   x + w * 0.88, y + h * 0.72, color=GOLD)

    elif kind == "shooting_line":
        soccer_field(slide, x, y, w, h)
        goal(slide, cx, y + Inches(0.3), width_in=1.4, orientation="horizontal")
        # Line of players
        line_y = y + h * 0.78
        for i in range(4):
            px = x + w * (0.2 + i * 0.12)
            player(slide, px, line_y, GOLD, "P")
            ball(slide, px, line_y + Inches(0.18))
        # Active shooter
        sx = x + w * 0.55
        sy = y + h * 0.5
        player(slide, sx, sy, GOLD, "P")
        ball(slide, sx + Inches(0.15), sy)
        run_arrow(slide, sx + Inches(0.18), sy,
                   cx, y + Inches(0.4), color=RED)

    elif kind == "finishing_pass":
        soccer_field(slide, x, y, w, h, penalty_boxes=True)
        goal(slide, cx, y + Inches(0.3), width_in=1.6, orientation="horizontal")
        coach = (x + w * 0.18, y + h * 0.55)
        runner_start = (x + w * 0.7, y + h * 0.78)
        finish = (cx + Inches(0.2), y + Inches(0.55))
        player(slide, coach[0], coach[1], GREEN, "C")
        player(slide, runner_start[0], runner_start[1], GOLD, "P")
        pass_arrow(slide, coach[0] + Inches(0.15), coach[1],
                    x + w * 0.55, y + h * 0.45)
        run_arrow(slide, runner_start[0] - Inches(0.1), runner_start[1] - Inches(0.05),
                   finish[0], finish[1] + Inches(0.05), color=GOLD)
        ball(slide, coach[0] + Inches(0.2), coach[1])

    elif kind == "mini_game":
        soccer_field(slide, x, y, w, h, halves=True)
        # Small goals at each end
        goal(slide, x + Inches(0.2), cy, width_in=0.7, orientation="vertical")
        goal(slide, x + w - Inches(0.2), cy, width_in=0.7, orientation="vertical")
        # Teams
        for px, py in [(0.28, 0.32), (0.34, 0.58), (0.4, 0.78)]:
            player(slide, x + w * px, y + h * py, GOLD, "A")
        for px, py in [(0.72, 0.32), (0.66, 0.58), (0.6, 0.78)]:
            player(slide, x + w * px, y + h * py, BLUE, "B")
        ball(slide, cx, cy)

    elif kind == "scrimmage":
        soccer_field(slide, x, y, w, h, halves=True, center_circle=True, penalty_boxes=True)
        goal(slide, x + Inches(0.18), cy, width_in=1.2, orientation="vertical")
        goal(slide, x + w - Inches(0.18), cy, width_in=1.2, orientation="vertical")
        # Team A
        for px, py in [(0.22, 0.5), (0.32, 0.3), (0.32, 0.7), (0.42, 0.5)]:
            player(slide, x + w * px, y + h * py, GOLD, "A")
        # Team B
        for px, py in [(0.78, 0.5), (0.68, 0.3), (0.68, 0.7), (0.58, 0.5)]:
            player(slide, x + w * px, y + h * py, BLUE, "B")
        # Ball at center
        ball(slide, cx, cy)
        # Sample passing line
        pass_arrow(slide, x + w * 0.42 + Inches(0.15), y + h * 0.5,
                    x + w * 0.55, y + h * 0.6)

    elif kind == "formation":
        soccer_field(slide, x, y, w, h, halves=True, center_circle=True, penalty_boxes=True)
        goal(slide, x + Inches(0.18), cy, width_in=1.2, orientation="vertical")
        goal(slide, x + w - Inches(0.18), cy, width_in=1.2, orientation="vertical")
        # 4-3-1 ish formation
        formation = [
            (0.22, 0.5, "GK", PURPLE),
            (0.32, 0.25, "D", BLUE),
            (0.32, 0.5, "D", BLUE),
            (0.32, 0.75, "D", BLUE),
            (0.44, 0.35, "M", GOLD),
            (0.44, 0.65, "M", GOLD),
            (0.55, 0.5, "F", RED),
        ]
        for px, py, lab, col in formation:
            player(slide, x + w * px, y + h * py, col, lab)
        ball(slide, x + w * 0.44, y + h * 0.5)

    else:
        grid_field(slide, x, y, w, h)
        player(slide, cx, cy, GOLD, "P")
        ball(slide, cx + Inches(0.2), cy)

    # Diagram caption strip — space label at the bottom of the diagram
    cap_h = Inches(0.32)
    cap = rect(slide, x, y + h - cap_h, w, cap_h, fill=BG_DARK)
    cap.fill.transparency = 0.4 if False else None
    textbox(slide, x + Inches(0.15), y + h - cap_h, w - Inches(0.3), cap_h,
            "📏  " + space, size=10, bold=True, color=WHITE_SOFT,
            anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT, letter_spacing=50)
    textbox(slide, x + Inches(0.15), y + h - cap_h, w - Inches(0.3), cap_h,
            tag.upper(), size=9, bold=True, color=band,
            anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT, letter_spacing=80)


# ── Slide builders ───────────────────────────────────────────────────────────

def add_cover_slide(prs, week, plan_meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height

    # Dark gradient background
    bg = rect(slide, 0, 0, sw, sh)
    fill_gradient(bg, BG_DARK, RGBColor(0x07, 0x14, 0x0E), 135.0)

    # Big diagonal green accent panel on the right
    accent_w = sw // 2 + Inches(0.5)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     sw - accent_w + Inches(2.5), -Inches(1),
                                     accent_w, sh + Inches(2))
    fill_gradient(accent, GREEN_DEEP, GREEN, 135.0)
    no_line(accent)
    accent.rotation = 8

    # Soft circle decoration
    deco = oval(slide, sw - Inches(3.5), sh - Inches(3.5), Inches(3.5), Inches(3.5),
                fill=GREEN_BRIGHT)
    deco.fill.transparency = None

    # Top accent line
    rect(slide, 0, 0, sw, Inches(0.1), fill=GREEN_BRIGHT)

    # qmtk logo mark (text-based)
    textbox(slide, Inches(0.7), Inches(0.55), Inches(4), Inches(0.5),
            "qmtk Soccer", size=20, bold=True, color=GREEN_BRIGHT, letter_spacing=120)
    textbox(slide, Inches(0.7), Inches(0.92), Inches(4), Inches(0.35),
            "COACH PLAYBOOK", size=10, bold=True, color=TEXT_MID, letter_spacing=400)

    # Title block
    d = datetime.strptime(week["weekOf"], "%Y-%m-%d")
    week_str = f"Week of {MONTHS_FULL[d.month - 1]} {d.day}"
    textbox(slide, Inches(0.7), Inches(2.6), Inches(8.5), Inches(0.5),
            "TRAINING SESSION", size=14, bold=True, color=GREEN_BRIGHT, letter_spacing=400)
    textbox(slide, Inches(0.7), Inches(3.05), Inches(9.5), Inches(1.4),
            week_str, size=60, bold=True, color=WHITE)
    textbox(slide, Inches(0.7), Inches(4.5), Inches(9.5), Inches(0.6),
            week.get("label", ""), size=22, color=GREEN_BRIGHT, italic=True)

    # Info strip
    strip_y = Inches(5.6)
    rect(slide, Inches(0.7), strip_y, Inches(0.1), Inches(1.1), fill=GREEN_BRIGHT)
    textbox(slide, Inches(0.95), strip_y, Inches(5), Inches(0.35),
            "LOCATION", size=9, bold=True, color=TEXT_DIM, letter_spacing=300)
    textbox(slide, Inches(0.95), strip_y + Inches(0.35), Inches(8), Inches(0.4),
            plan_meta.get("location", ""), size=14, bold=True, color=WHITE)
    textbox(slide, Inches(0.95), strip_y + Inches(0.75), Inches(8), Inches(0.4),
            plan_meta.get("time", ""), size=12, color=TEXT_MID)

    # Footer
    textbox(slide, Inches(0.7), sh - Inches(0.55), Inches(10), Inches(0.3),
            "qmtk.org  ·  Train Smart. Play Smart.",
            size=10, bold=True, color=TEXT_DIM, letter_spacing=250)


def add_section_divider(prs, big_number, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height

    bg = rect(slide, 0, 0, sw, sh)
    fill_gradient(bg, GREEN_DEEP, GREEN, 135.0)

    # Big number in the corner, watermark style
    textbox(slide, Inches(7.5), Inches(0.5), Inches(5.5), Inches(6),
            big_number, size=320, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    # Make it watermark-y via transparency (PowerPoint will keep solid, but visual is bold)

    rect(slide, 0, 0, sw, Inches(0.12), fill=GREEN_BRIGHT)
    rect(slide, Inches(0.8), Inches(3.0), Inches(0.7), Inches(0.05), fill=GREEN_BRIGHT)
    textbox(slide, Inches(0.8), Inches(3.15), Inches(10), Inches(0.5),
            "SECTION", size=12, bold=True, color=GREEN_BRIGHT, letter_spacing=500)
    textbox(slide, Inches(0.8), Inches(3.55), Inches(11), Inches(1.4),
            title, size=56, bold=True, color=WHITE)
    textbox(slide, Inches(0.8), Inches(5.05), Inches(11), Inches(0.5),
            subtitle, size=20, color=RGBColor(0xC5, 0xE8, 0xCE), italic=True)

    textbox(slide, Inches(0.8), sh - Inches(0.5), Inches(10), Inches(0.3),
            "qmtk Soccer · Coach Playbook", size=10, bold=True,
            color=RGBColor(0xC5, 0xE8, 0xCE), letter_spacing=200)


def add_drill_slide(prs, age_label, step, index_str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height

    # Background
    bg = rect(slide, 0, 0, sw, sh)
    fill_gradient(bg, BG_DARK, BG_PANEL, 135.0)

    band_color = TAG_COLOR.get(step["tag"], GREEN)

    # Top colored accent strip
    rect(slide, 0, 0, sw, Inches(0.16), fill=band_color)

    # Brand top-left
    textbox(slide, Inches(0.5), Inches(0.32), Inches(4), Inches(0.3),
            "qmtk Soccer", size=11, bold=True, color=GREEN_BRIGHT, letter_spacing=200)
    textbox(slide, Inches(0.5), Inches(0.58), Inches(5), Inches(0.3),
            "COACH PLAYBOOK · " + age_label.upper(), size=8, bold=True,
            color=TEXT_DIM, letter_spacing=400)

    # Index top-right
    textbox(slide, sw - Inches(1.5), Inches(0.32), Inches(1.1), Inches(0.45),
            index_str, size=14, bold=True, color=TEXT_MID, align=PP_ALIGN.RIGHT)

    # Title block
    chip_x = Inches(0.5)
    chip_y = Inches(1.2)
    chip_shape, chip_w = chip(slide, chip_x, chip_y, step["tag"].upper(),
                               bg=band_color, fg=BG_DARK, size=11)

    # Time chip next to tag
    time_chip, time_w = chip(slide, chip_x + chip_w + Inches(0.12), chip_y,
                              step["time"], bg=BG_CARD, fg=WHITE_SOFT, size=10)

    # Space chip next to time
    space_label = "📏  " + step["space"]
    chip(slide, chip_x + chip_w + time_w + Inches(0.24), chip_y,
         space_label, bg=BG_CARD, fg=WHITE_SOFT, size=10)

    textbox(slide, Inches(0.5), Inches(1.65), Inches(8.5), Inches(1.1),
            step["title"], size=36, bold=True, color=WHITE)

    # Underline accent
    rect(slide, Inches(0.5), Inches(2.55), Inches(0.7), Inches(0.06), fill=band_color)

    # Left column: description + cues
    left_x = Inches(0.5)
    left_w = Inches(5.6)
    body_y = Inches(2.85)

    textbox(slide, left_x, body_y, left_w, Inches(0.35),
            "THE DRILL", size=10, bold=True, color=band_color, letter_spacing=400)
    textbox(slide, left_x, body_y + Inches(0.35), left_w, Inches(1.9),
            step["desc"], size=13, color=WHITE_SOFT)

    cue_label_y = body_y + Inches(2.4)
    textbox(slide, left_x, cue_label_y, left_w, Inches(0.35),
            "KEY COACHING CUES", size=10, bold=True, color=band_color, letter_spacing=400)

    cue_y = cue_label_y + Inches(0.4)
    # Cue card
    cue_h = Inches(2.0)
    cue_card = rect(slide, left_x, cue_y, left_w, cue_h, fill=BG_CARD, rounded=True,
                     line=BORDER, line_w=1.0)
    # Cue items
    cue_inner_y = cue_y + Inches(0.2)
    for cue in step.get("focus", []):
        # Bullet dot
        dot_r = Inches(0.06)
        oval(slide, left_x + Inches(0.25), cue_inner_y + Inches(0.08),
             dot_r * 2, dot_r * 2, fill=band_color)
        textbox(slide, left_x + Inches(0.6), cue_inner_y, left_w - Inches(0.8), Inches(0.45),
                cue, size=12, color=WHITE_SOFT)
        cue_inner_y += Inches(0.46)

    # Right column: diagram
    diag_x = Inches(6.4)
    diag_y = Inches(2.85)
    diag_w = Inches(6.4)
    diag_h = Inches(3.95)

    # Diagram backing card
    backing = rect(slide, diag_x - Inches(0.15), diag_y - Inches(0.15),
                    diag_w + Inches(0.3), diag_h + Inches(0.3),
                    fill=BG_CARD, rounded=True, line=BORDER, line_w=1.0, shadow=True)

    draw_diagram(slide, diag_x, diag_y, diag_w, diag_h,
                 title=step["title"], tag=step["tag"], space=step["space"])

    # Legend below diagram
    leg_y = diag_y + diag_h + Inches(0.3)
    leg_items = [
        (GOLD, "Player A / Attacker"),
        (BLUE, "Player B / Receiver"),
        (PINK_RED, "Defender"),
        (WHITE, "Ball"),
    ]
    cx = diag_x
    for col, lbl in leg_items:
        r = Inches(0.09)
        oval(slide, cx, leg_y + Inches(0.06), r * 2, r * 2, fill=col, line=BG_DARK, line_w=0.75)
        textbox(slide, cx + Inches(0.22), leg_y, Inches(1.5), Inches(0.3),
                lbl, size=9, color=TEXT_MID)
        cx += Inches(1.5)

    # Footer
    foot_y = sh - Inches(0.35)
    rect(slide, 0, foot_y, sw, Inches(0.35), fill=BG_PANEL)
    textbox(slide, Inches(0.5), foot_y, Inches(8), Inches(0.35),
            f"qmtk Soccer · Coach Playbook · {age_label}",
            size=9, color=TEXT_DIM, anchor=MSO_ANCHOR.MIDDLE, letter_spacing=200)
    textbox(slide, sw - Inches(2.5), foot_y, Inches(2), Inches(0.35),
            "qmtk.org", size=9, color=TEXT_DIM, align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE, letter_spacing=200)


# ── Deck assembly ────────────────────────────────────────────────────────────

def build_week(plan_meta, week):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover_slide(prs, week, plan_meta)

    prek = week.get("prek", [])
    elem = week.get("elem", [])

    add_section_divider(prs, "01", "Pre-K & Kindergarten", "Ages 3–5  ·  45 minute session")
    for i, step in enumerate(prek):
        add_drill_slide(prs, "Pre-K & K", step, f"01.{i + 1}  /  04")

    add_section_divider(prs, "02", "Grades 1–5", "Ages 6–11  ·  60 minute session")
    for i, step in enumerate(elem):
        add_drill_slide(prs, "Grades 1–5", step, f"02.{i + 1}  /  04")

    return prs


def main():
    with open(JSON_PATH, "r", encoding="utf-8") as f:
        data = json.load(f)

    os.makedirs(OUT_DIR, exist_ok=True)
    plan_meta = {"location": data.get("location", ""), "time": data.get("time", "")}

    for week in data.get("weeks", []):
        prs = build_week(plan_meta, week)
        out = os.path.join(OUT_DIR, f"{week['weekOf']}.pptx")
        prs.save(out)
        print(f"wrote {out}")


if __name__ == "__main__":
    main()
